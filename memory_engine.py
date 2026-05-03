import os
import pypdf
from docx import Document
from pptx import Presentation
import logging

class MemoryEngine:
    def __init__(self, vault_path: str):
        self.vault_path = vault_path
        self.knowledge_summary = ""
        
    def extract_text_from_pdf(self, file_path: str, max_pages: int = 5) -> str:
        text = ""
        try:
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f)
                for i in range(min(len(reader.pages), max_pages)):
                    page_text = reader.pages[i].extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            logging.error(f"Error reading PDF {file_path}: {e}")
        return text

    def extract_text_from_docx(self, file_path: str, max_paragraphs: int = 50) -> str:
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs[:max_paragraphs]:
                text += paragraph.text + "\n"
        except Exception as e:
            logging.error(f"Error reading DOCX {file_path}: {e}")
        return text

    def extract_text_from_pptx(self, file_path: str) -> str:
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides[:10]: # First 10 slides
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            logging.error(f"Error reading PPTX {file_path}: {e}")
        return text

    def refresh_knowledge(self):
        """Scans the vault and builds a dense summary of available knowledge."""
        summary_parts = []
        if not os.path.exists(self.vault_path):
            return "Vault path not found."
            
        files = os.listdir(self.vault_path)
        summary_parts.append(f"Found {len(files)} documents in Energy Vault.")
        
        # Focus on key documents for the summary
        key_files = [f for f in files if any(k in f.lower() for k in [
            'cv', 'resume', 'sop', 'energetika', 'mashina', 'crypto', 
            'gidravlik', 'hydraulic', 'contact', 'jamshid', 'olimov'
        ])]
        
        for file_name in key_files[:15]: # Limit to top 15 key files for prompt context
            path = os.path.join(self.vault_path, file_name)
            content = ""
            if file_name.endswith('.pdf'):
                content = self.extract_text_from_pdf(path)
            elif file_name.endswith('.docx'):
                content = self.extract_text_from_docx(path)
            elif file_name.endswith('.pptx'):
                content = self.extract_text_from_pptx(path)
            elif file_name.endswith('.txt'):
                try:
                    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read(5000)
                except:
                    pass
                
            if content:
                # Truncate content for the summary
                summary_parts.append(f"--- DOCUMENT: {file_name} ---\n{content[:800]}...\n")
        
        self.knowledge_summary = "\n".join(summary_parts)
        return self.knowledge_summary

    def get_context(self) -> str:
        if not self.knowledge_summary:
            self.refresh_knowledge()
        return self.knowledge_summary
